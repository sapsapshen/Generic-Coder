import os, sys, json, io, base64, tempfile, subprocess, hashlib
from enum import Enum, auto
from pathlib import Path
from typing import Optional, Dict, List, Any, Union
from datetime import datetime

class MediaType(Enum):
    IMAGE = auto()
    VIDEO = auto()
    AUDIO = auto()
    PDF = auto()
    DOCUMENT = auto()
    SPREADSHEET = auto()
    PRESENTATION = auto()
    ARCHIVE = auto()
    CODE = auto()
    DATA = auto()
    UNKNOWN = auto()

MIME_TO_TYPE = {
    'image/png': MediaType.IMAGE, 'image/jpeg': MediaType.IMAGE,
    'image/gif': MediaType.IMAGE, 'image/webp': MediaType.IMAGE,
    'image/svg+xml': MediaType.IMAGE, 'image/bmp': MediaType.IMAGE,
    'image/tiff': MediaType.IMAGE, 'image/x-icon': MediaType.IMAGE,
    'video/mp4': MediaType.VIDEO, 'video/quicktime': MediaType.VIDEO,
    'video/x-msvideo': MediaType.VIDEO, 'video/webm': MediaType.VIDEO,
    'video/x-matroska': MediaType.VIDEO, 'video/x-flv': MediaType.VIDEO,
    'audio/mpeg': MediaType.AUDIO, 'audio/wav': MediaType.AUDIO,
    'audio/ogg': MediaType.AUDIO, 'audio/flac': MediaType.AUDIO,
    'audio/aac': MediaType.AUDIO, 'audio/x-m4a': MediaType.AUDIO,
    'application/pdf': MediaType.PDF,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': MediaType.DOCUMENT,
    'application/msword': MediaType.DOCUMENT,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': MediaType.SPREADSHEET,
    'application/vnd.ms-excel': MediaType.SPREADSHEET,
    'text/csv': MediaType.DATA,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': MediaType.PRESENTATION,
    'application/zip': MediaType.ARCHIVE, 'application/x-tar': MediaType.ARCHIVE,
    'application/gzip': MediaType.ARCHIVE, 'application/x-7z-compressed': MediaType.ARCHIVE,
    'application/x-rar-compressed': MediaType.ARCHIVE, 'application/x-bzip2': MediaType.ARCHIVE,
    'application/json': MediaType.DATA, 'text/xml': MediaType.DATA,
    'application/xml': MediaType.DATA, 'text/yaml': MediaType.DATA,
    'text/markdown': MediaType.CODE, 'text/html': MediaType.CODE,
}

EXT_TO_TYPE = {
    '.png': MediaType.IMAGE, '.jpg': MediaType.IMAGE, '.jpeg': MediaType.IMAGE,
    '.gif': MediaType.IMAGE, '.webp': MediaType.IMAGE, '.svg': MediaType.IMAGE,
    '.bmp': MediaType.IMAGE, '.tiff': MediaType.IMAGE, '.ico': MediaType.IMAGE,
    '.heic': MediaType.IMAGE, '.heif': MediaType.IMAGE,
    '.mp4': MediaType.VIDEO, '.mov': MediaType.VIDEO, '.avi': MediaType.VIDEO,
    '.webm': MediaType.VIDEO, '.mkv': MediaType.VIDEO, '.flv': MediaType.VIDEO,
    '.wmv': MediaType.VIDEO, '.m4v': MediaType.VIDEO,
    '.mp3': MediaType.AUDIO, '.wav': MediaType.AUDIO, '.ogg': MediaType.AUDIO,
    '.flac': MediaType.AUDIO, '.aac': MediaType.AUDIO, '.m4a': MediaType.AUDIO,
    '.wma': MediaType.AUDIO, '.opus': MediaType.AUDIO,
    '.pdf': MediaType.PDF,
    '.docx': MediaType.DOCUMENT, '.doc': MediaType.DOCUMENT,
    '.xlsx': MediaType.SPREADSHEET, '.xls': MediaType.SPREADSHEET, '.csv': MediaType.DATA,
    '.pptx': MediaType.PRESENTATION, '.ppt': MediaType.PRESENTATION,
    '.zip': MediaType.ARCHIVE, '.tar': MediaType.ARCHIVE, '.gz': MediaType.ARCHIVE,
    '.7z': MediaType.ARCHIVE, '.rar': MediaType.ARCHIVE, '.bz2': MediaType.ARCHIVE,
    '.xz': MediaType.ARCHIVE,
    '.json': MediaType.DATA, '.xml': MediaType.DATA, '.yaml': MediaType.DATA,
    '.yml': MediaType.DATA, '.toml': MediaType.DATA,
    '.py': MediaType.CODE, '.js': MediaType.CODE, '.ts': MediaType.CODE,
    '.java': MediaType.CODE, '.c': MediaType.CODE, '.cpp': MediaType.CODE,
    '.h': MediaType.CODE, '.rs': MediaType.CODE, '.go': MediaType.CODE,
    '.rb': MediaType.CODE, '.sh': MediaType.CODE, '.bash': MediaType.CODE,
    '.html': MediaType.CODE, '.css': MediaType.CODE, '.scss': MediaType.CODE,
    '.sql': MediaType.CODE, '.swift': MediaType.CODE, '.kt': MediaType.CODE,
    '.md': MediaType.CODE, '.txt': MediaType.CODE,
}

try:
    from PIL import Image
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

try:
    import PyPDF2
    _PYPDF_AVAILABLE = True
except ImportError:
    _PYPDF_AVAILABLE = False

try:
    import pdfplumber
    _PDFPLUMBER_AVAILABLE = True
except ImportError:
    _PDFPLUMBER_AVAILABLE = False

def _get_ffprobe_path():
    for p in ['ffprobe', '/usr/local/bin/ffprobe', '/opt/homebrew/bin/ffprobe']:
        try:
            subprocess.run([p, '-version'], capture_output=True, timeout=3)
            return p
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None

def _get_ffmpeg_path():
    for p in ['ffmpeg', '/usr/local/bin/ffmpeg', '/opt/homebrew/bin/ffmpeg']:
        try:
            subprocess.run([p, '-version'], capture_output=True, timeout=3)
            return p
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None

_FFPROBE = _get_ffprobe_path()
_FFMPEG = _get_ffmpeg_path()


class MediaHandler:
    def __init__(self):
        self._preview_cache_dir = Path(tempfile.gettempdir()) / 'genericagent_previews'
        self._preview_cache_dir.mkdir(parents=True, exist_ok=True)

    def detect_type(self, file_path: str) -> MediaType:
        ext = os.path.splitext(file_path)[1].lower()
        if ext in EXT_TO_TYPE:
            return EXT_TO_TYPE[ext]
        return MediaType.UNKNOWN

    def get_file_info(self, file_path: str) -> dict:
        if not os.path.exists(file_path):
            return {'status': 'error', 'msg': f'File not found: {file_path}'}

        path = Path(file_path)
        stat = path.stat()
        info = {
            'status': 'success',
            'name': path.name,
            'path': str(path.absolute()),
            'size': stat.st_size,
            'size_human': self._format_size(stat.st_size),
            'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'extension': path.suffix.lower(),
            'media_type': self.detect_type(file_path).name,
        }

        media_type = self.detect_type(file_path)

        if media_type == MediaType.IMAGE and _PIL_AVAILABLE:
            try:
                with Image.open(file_path) as img:
                    info['image'] = {
                        'width': img.width,
                        'height': img.height,
                        'format': img.format,
                        'mode': img.mode,
                        'aspect_ratio': round(img.width / img.height, 3) if img.height else 0,
                    }
            except Exception:
                pass

        if media_type == MediaType.PDF:
            info['pdf'] = self._get_pdf_info(file_path)

        if media_type == MediaType.VIDEO or media_type == MediaType.AUDIO:
            probe = self._ffprobe(file_path)
            if probe:
                info['media'] = probe

        return info

    def _format_size(self, size: int) -> str:
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if size < 1024:
                return f'{size:.1f} {unit}'
            size /= 1024
        return f'{size:.1f} PB'

    def _get_pdf_info(self, file_path: str) -> dict:
        info = {}
        if _PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(file_path) as pdf:
                    info['pages'] = len(pdf.pages)
                    if pdf.metadata:
                        info['metadata'] = {
                            k: str(v) for k, v in (pdf.metadata or {}).items() if v
                        }
            except Exception:
                pass
        if not info and _PYPDF_AVAILABLE:
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    info['pages'] = len(reader.pages)
                    if reader.metadata:
                        info['metadata'] = {
                            k[1:] if k.startswith('/') else k: str(v)
                            for k, v in (reader.metadata or {}).items() if v
                        }
            except Exception:
                pass
        if not info:
            try:
                result = subprocess.run(
                    ['python3', '-c', f"""
import fitz
doc = fitz.open('{file_path}')
print(len(doc))
"""], capture_output=True, text=True, timeout=10)
                if result.returncode == 0:
                    info['pages'] = int(result.stdout.strip())
            except Exception:
                pass
        return info

    def _ffprobe(self, file_path: str) -> Optional[dict]:
        if not _FFPROBE:
            return None
        try:
            result = subprocess.run(
                [_FFPROBE, '-v', 'quiet', '-print_format', 'json',
                 '-show_format', '-show_streams', file_path],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode != 0:
                return None
            data = json.loads(result.stdout)

            info = {}
            fmt = data.get('format', {})
            info['duration'] = float(fmt.get('duration', 0))
            info['duration_human'] = self._format_duration(info['duration'])
            info['bit_rate'] = int(fmt.get('bit_rate', 0))
            info['size'] = int(fmt.get('size', 0))
            info['format_name'] = fmt.get('format_name', '')

            for stream in data.get('streams', []):
                codec_type = stream.get('codec_type', '')
                if codec_type == 'video':
                    info['video'] = {
                        'codec': stream.get('codec_name', ''),
                        'width': stream.get('width', 0),
                        'height': stream.get('height', 0),
                        'fps': self._parse_fps(stream.get('r_frame_rate', '')),
                        'aspect_ratio': stream.get('display_aspect_ratio', ''),
                    }
                elif codec_type == 'audio':
                    info['audio'] = {
                        'codec': stream.get('codec_name', ''),
                        'channels': stream.get('channels', 0),
                        'sample_rate': int(stream.get('sample_rate', 0)),
                        'bit_rate': int(stream.get('bit_rate', 0)),
                    }
            return info
        except Exception:
            return None

    def _parse_fps(self, fps_str: str) -> float:
        if not fps_str:
            return 0.0
        try:
            num, den = fps_str.split('/')
            den = int(den)
            if den == 0:
                return 0.0
            return round(int(num) / den, 2)
        except (ValueError, ZeroDivisionError):
            return 0.0

    def _format_duration(self, seconds: float) -> str:
        if not seconds or seconds <= 0:
            return '0:00'
        h = int(seconds // 3600)
        m = int((seconds % 3600) // 60)
        s = int(seconds % 60)
        if h > 0:
            return f'{h}:{m:02d}:{s:02d}'
        return f'{m}:{s:02d}'

    def generate_thumbnail(self, file_path: str, max_size: int = 400) -> Optional[str]:
        media_type = self.detect_type(file_path)
        file_hash = hashlib.md5(file_path.encode()).hexdigest()[:12]
        thumb_path = self._preview_cache_dir / f'{file_hash}_thumb.png'

        if thumb_path.exists():
            return str(thumb_path)

        if media_type == MediaType.IMAGE and _PIL_AVAILABLE:
            try:
                with Image.open(file_path) as img:
                    img.thumbnail((max_size, max_size), Image.LANCZOS)
                    img.save(thumb_path, 'PNG')
                return str(thumb_path)
            except Exception:
                return None

        if media_type == MediaType.VIDEO and _FFMPEG:
            try:
                subprocess.run(
                    [_FFMPEG, '-y', '-i', file_path, '-vframes', '1',
                     '-vf', f'scale={max_size}:{max_size}:force_original_aspect_ratio=decrease',
                     '-f', 'image2', str(thumb_path)],
                    capture_output=True, timeout=30
                )
                if thumb_path.exists():
                    return str(thumb_path)
            except Exception:
                pass

        if media_type == MediaType.PDF and _PIL_AVAILABLE:
            try:
                if _PDFPLUMBER_AVAILABLE:
                    with pdfplumber.open(file_path) as pdf:
                        if pdf.pages:
                            page = pdf.pages[0]
                            img = page.to_image(resolution=72)
                            img.save(str(thumb_path))
                            return str(thumb_path)
            except Exception:
                pass
            try:
                from PIL import Image
                subprocess.run(
                    ['python3', '-c', f"""
import fitz
doc = fitz.open('{file_path}')
page = doc[0]
pix = page.get_pixmap(dpi=72)
pix.save('{thumb_path}')
"""], capture_output=True, timeout=15)
                if thumb_path.exists():
                    return str(thumb_path)
            except Exception:
                pass

        return None

    def generate_preview_url(self, file_path: str) -> dict:
        thumb = self.generate_thumbnail(file_path)
        if not thumb:
            return {'status': 'error', 'msg': 'Could not generate preview'}
        try:
            with open(thumb, 'rb') as f:
                data = base64.b64encode(f.read()).decode()
            return {
                'status': 'success',
                'thumbnail': f'data:image/png;base64,{data}',
                'path': thumb
            }
        except Exception as e:
            return {'status': 'error', 'msg': str(e)}

    def extract_text(self, file_path: str) -> dict:
        media_type = self.detect_type(file_path)
        text = ''

        if media_type == MediaType.PDF:
            text = self._extract_pdf_text(file_path)
        elif media_type == MediaType.DOCUMENT:
            text = self._extract_docx_text(file_path)
        elif media_type == MediaType.SPREADSHEET:
            text = self._extract_xlsx_text(file_path)
        elif media_type == MediaType.PRESENTATION:
            text = self._extract_pptx_text(file_path)
        elif media_type in (MediaType.CODE, MediaType.DATA):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read()
            except Exception:
                pass
        else:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read()
            except Exception:
                pass

        if text:
            return {
                'status': 'success',
                'text': text,
                'length': len(text),
                'preview': text[:500]
            }
        return {'status': 'error', 'msg': 'Could not extract text from this file type'}

    def _extract_pdf_text(self, file_path: str) -> str:
        pages = []
        if _PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text()
                        if t:
                            pages.append(t)
            except Exception:
                pass
        if not pages and _PYPDF_AVAILABLE:
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        t = page.extract_text()
                        if t:
                            pages.append(t)
            except Exception:
                pass
        return '\n\n'.join(pages)

    def _extract_docx_text(self, file_path: str) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs if p.text)
        except ImportError:
            pass
        try:
            result = subprocess.run(
                ['python3', '-c', f"""
import zipfile, xml.etree.ElementTree as ET
z = zipfile.ZipFile('{file_path}')
xml_content = z.read('word/document.xml')
tree = ET.fromstring(xml_content)
ns = {{'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}}
texts = [node.text for node in tree.iter('{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}t') if node.text]
print('\\n'.join(texts))
"""], capture_output=True, text=True, timeout=15)
            if result.returncode == 0:
                return result.stdout.strip()
        except Exception:
            pass
        return ''

    def _extract_xlsx_text(self, file_path: str) -> str:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            rows = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows.append(f'--- Sheet: {sheet_name} ---')
                for row in ws.iter_rows(values_only=True):
                    rows.append('\t'.join(str(c) if c is not None else '' for c in row))
                    if len(rows) > 1000:
                        rows.append('... (truncated)')
                        break
            wb.close()
            return '\n'.join(rows)
        except ImportError:
            pass
        if file_path.endswith('.csv'):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
            except Exception:
                pass
        return ''

    def _extract_pptx_text(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        texts.append(shape.text)
            return '\n\n'.join(texts)
        except ImportError:
            pass
        return ''

    def convert_image(self, file_path: str, target_format: str,
                     quality: int = 85, max_size: int = 0) -> Optional[str]:
        if not _PIL_AVAILABLE:
            return None
        target_format = target_format.lower().lstrip('.')
        out_path = str(Path(file_path).with_suffix(f'.{target_format}'))
        try:
            with Image.open(file_path) as img:
                if img.mode in ('RGBA', 'P') and target_format in ('jpg', 'jpeg'):
                    img = img.convert('RGB')
                if max_size:
                    img.thumbnail((max_size, max_size), Image.LANCZOS)
                img.save(out_path, format=target_format.upper(), quality=quality)
            return out_path
        except Exception:
            return None

    def extract_audio(self, video_path: str, output_format: str = 'mp3') -> Optional[str]:
        if not _FFMPEG:
            return None
        out_path = str(Path(video_path).with_suffix(f'.{output_format}'))
        try:
            subprocess.run(
                [_FFMPEG, '-y', '-i', video_path, '-vn', '-acodec',
                 'libmp3lame' if output_format == 'mp3' else 'copy',
                 out_path],
                capture_output=True, timeout=120
            )
            if os.path.exists(out_path):
                return out_path
        except Exception:
            pass
        return None

    def extract_video_frame(self, video_path: str, time_seconds: float = 0,
                           output_path: str = '') -> Optional[str]:
        if not _FFMPEG:
            return None
        if not output_path:
            output_path = str(Path(video_path).with_suffix('.frame.png'))
        try:
            subprocess.run(
                [_FFMPEG, '-y', '-ss', str(time_seconds), '-i', video_path,
                 '-vframes', '1', '-q:v', '2', output_path],
                capture_output=True, timeout=30
            )
            if os.path.exists(output_path):
                return output_path
        except Exception:
            pass
        return None

    def convert_video(self, video_path: str, target_format: str = 'mp4',
                     crf: int = 23, max_width: int = 0) -> Optional[str]:
        if not _FFMPEG:
            return None
        out_path = str(Path(video_path).with_suffix(f'.{target_format}'))
        try:
            cmd = [_FFMPEG, '-y', '-i', video_path, '-crf', str(crf)]
            if max_width:
                cmd += ['-vf', f'scale={max_width}:-1']
            cmd.append(out_path)
            subprocess.run(cmd, capture_output=True, timeout=300)
            if os.path.exists(out_path):
                return out_path
        except Exception:
            pass
        return None

    def extract_archive(self, archive_path: str, output_dir: str = '') -> dict:
        if not output_dir:
            output_dir = str(Path(archive_path).with_suffix(''))
        os.makedirs(output_dir, exist_ok=True)

        import zipfile, tarfile
        ext = os.path.splitext(archive_path)[1].lower()

        try:
            if ext == '.zip':
                with zipfile.ZipFile(archive_path, 'r') as zf:
                    zf.extractall(output_dir)
                    files = zf.namelist()
            elif ext in ('.tar', '.gz', '.bz2', '.xz'):
                mode_map = {'.tar': 'r', '.gz': 'r:gz', '.bz2': 'r:bz2', '.xz': 'r:xz'}
                mode = mode_map.get(ext, 'r:*')
                with tarfile.open(archive_path, mode) as tf:
                    tf.extractall(output_dir)
                    files = tf.getnames()
            else:
                return {'status': 'error', 'msg': f'Unsupported archive format: {ext}'}
            return {'status': 'success', 'output_dir': output_dir, 'files': files[:200]}
        except Exception as e:
            return {'status': 'error', 'msg': str(e)}

    def list_archive(self, archive_path: str) -> dict:
        import zipfile, tarfile
        ext = os.path.splitext(archive_path)[1].lower()
        try:
            if ext == '.zip':
                with zipfile.ZipFile(archive_path, 'r') as zf:
                    files = [{'name': n, 'size': i.file_size, 'is_dir': i.is_dir()}
                            for n, i in zip(zf.namelist(), zf.infolist())]
            elif ext in ('.tar', '.gz', '.bz2', '.xz'):
                mode_map = {'.tar': 'r', '.gz': 'r:gz', '.bz2': 'r:bz2', '.xz': 'r:xz'}
                mode = mode_map.get(ext, 'r:*')
                with tarfile.open(archive_path, mode) as tf:
                    files = [{'name': m.name, 'size': m.size, 'is_dir': m.isdir()}
                            for m in tf.getmembers()]
            else:
                return {'status': 'error', 'msg': f'Unsupported archive format: {ext}'}
            return {'status': 'success', 'files': files[:200], 'total': len(files)}
        except Exception as e:
            return {'status': 'error', 'msg': str(e)}


_global_media_handler: Optional[MediaHandler] = None

def get_media_handler() -> MediaHandler:
    global _global_media_handler
    if _global_media_handler is None:
        _global_media_handler = MediaHandler()
    return _global_media_handler